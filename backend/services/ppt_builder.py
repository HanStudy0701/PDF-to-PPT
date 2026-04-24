from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from backend.models.ir import Element, PresentationIR

EMU_PER_INCH = 914400
PT_PER_INCH = 72


def _hex_to_rgb(color: str | None) -> RGBColor:
    if not color:
        return RGBColor(0, 0, 0)
    c = color.lstrip("#")
    if len(c) != 6:
        return RGBColor(0, 0, 0)
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))


def _pt_to_inches(v: float) -> float:
    return v / PT_PER_INCH


def _add_text(slide, el: Element):
    box = slide.shapes.add_textbox(
        Inches(_pt_to_inches(el.x)),
        Inches(_pt_to_inches(el.y)),
        Inches(_pt_to_inches(el.width)),
        Inches(_pt_to_inches(el.height)),
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = el.text or ""
    if el.font_size:
        run.font.size = Pt(el.font_size)
    if el.font_family:
        run.font.name = el.font_family
    run.font.bold = bool(el.bold)
    run.font.italic = bool(el.italic)
    run.font.underline = bool(el.underline)
    run.font.color.rgb = _hex_to_rgb(el.font_color)


def _add_image(slide, el: Element):
    if not el.path or not Path(el.path).exists():
        return
    slide.shapes.add_picture(
        el.path,
        Inches(_pt_to_inches(el.x)),
        Inches(_pt_to_inches(el.y)),
        Inches(_pt_to_inches(el.width)),
        Inches(_pt_to_inches(el.height)),
    )


def _add_shape(slide, el: Element):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if (el.border_radius or 0) > 0 else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(_pt_to_inches(el.x)),
        Inches(_pt_to_inches(el.y)),
        Inches(_pt_to_inches(el.width)),
        Inches(_pt_to_inches(el.height)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(el.fill_color)
    shape.line.color.rgb = _hex_to_rgb(el.line_color)
    shape.line.width = Pt(el.line_width or 1)


def build_pptx(ir: PresentationIR, output_pptx: Path) -> Path:
    prs = Presentation()
    if ir.slides:
        first = ir.slides[0]
        prs.slide_width = Inches(_pt_to_inches(first.width))
        prs.slide_height = Inches(_pt_to_inches(first.height))

    blank_layout = prs.slide_layouts[6]

    for s in ir.slides:
        slide = prs.slides.add_slide(blank_layout)
        if s.background.path and Path(s.background.path).exists():
            slide.shapes.add_picture(
                s.background.path,
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

        for el in sorted(s.elements, key=lambda e: e.z_index):
            if el.type == "text":
                _add_text(slide, el)
            elif el.type in {"image", "icon"}:
                _add_image(slide, el)
            elif el.type in {"shape", "table_cell", "line"}:
                _add_shape(slide, el)

    prs.save(output_pptx)
    return output_pptx
